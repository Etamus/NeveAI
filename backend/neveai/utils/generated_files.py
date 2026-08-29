import csv
import io
import json
import re
import zipfile
from pathlib import Path, PurePosixPath
from typing import Any


MAX_SOURCE_CHARS = 4_000_000
MAX_OUTPUT_BYTES = 32 * 1024 * 1024
MAX_ARCHIVE_FILES = 100

TEXT_FORMATS = {
    "txt": "text/plain; charset=utf-8",
    "md": "text/markdown; charset=utf-8",
    "csv": "text/csv; charset=utf-8",
    "json": "application/json",
    "html": "text/html; charset=utf-8",
    "css": "text/css; charset=utf-8",
    "js": "text/javascript; charset=utf-8",
    "ts": "text/typescript; charset=utf-8",
    "py": "text/x-python; charset=utf-8",
    "java": "text/x-java-source; charset=utf-8",
    "c": "text/x-c; charset=utf-8",
    "cpp": "text/x-c++; charset=utf-8",
    "h": "text/x-c; charset=utf-8",
    "sh": "application/x-sh; charset=utf-8",
    "yaml": "application/yaml; charset=utf-8",
    "yml": "application/yaml; charset=utf-8",
    "xml": "application/xml; charset=utf-8",
    "sql": "application/sql; charset=utf-8",
    "rtf": "application/rtf",
}

MIME_TYPES = {
    **TEXT_FORMATS,
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "pdf": "application/pdf",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "zip": "application/zip",
}

FORMAT_ALIASES = {
    "markdown": "md",
    "javascript": "js",
    "typescript": "ts",
    "python": "py",
    "excel": "xlsx",
    "spreadsheet": "xlsx",
    "word": "docx",
    "powerpoint": "pptx",
    "presentation": "pptx",
}


class GeneratedFileError(ValueError):
    pass


def _normalize_format(value: str) -> str:
    normalized = (value or "").strip().lower().lstrip(".")
    return FORMAT_ALIASES.get(normalized, normalized)


def _sanitize_filename(filename: str, requested_format: str) -> tuple[str, str]:
    raw_name = Path((filename or "").strip()).name
    raw_name = re.sub(r"[\x00-\x1f<>:\"/\\|?*]", "_", raw_name).strip(" .")
    if not raw_name:
        raw_name = "arquivo"

    suffix_format = _normalize_format(Path(raw_name).suffix)
    file_format = _normalize_format(requested_format) or suffix_format or "txt"
    if file_format not in MIME_TYPES:
        raise GeneratedFileError(
            f"Formato '{file_format}' não suportado. Use: {', '.join(sorted(MIME_TYPES))}."
        )

    stem = Path(raw_name).stem if Path(raw_name).suffix else raw_name
    stem = stem[:120].strip(" .") or "arquivo"
    return f"{stem}.{file_format}", file_format


def _load_json(content: str, expected: str) -> Any:
    try:
        return json.loads(content)
    except json.JSONDecodeError as exc:
        raise GeneratedFileError(
            f"O conteúdo de {expected} deve ser JSON válido: {exc.msg}."
        ) from exc


def _strip_inline_markdown(text: str) -> str:
    text = re.sub(r"!\[([^\]]*)\]\([^)]*\)", r"\1", text)
    text = re.sub(r"\[([^\]]+)\]\(([^)]+)\)", r"\1 (\2)", text)
    text = re.sub(r"(`{1,3}|\*\*|__|~~)(.*?)(\1)", r"\2", text)
    return text.replace("\\|", "|").strip()


def _markdown_rows(lines: list[str], start: int) -> tuple[list[list[str]], int] | None:
    if start + 1 >= len(lines) or "|" not in lines[start]:
        return None
    separator = lines[start + 1].strip().strip("|")
    if not separator or not all(
        re.fullmatch(r":?-{3,}:?", cell.strip())
        for cell in separator.split("|")
    ):
        return None

    rows = []
    index = start
    while index < len(lines) and "|" in lines[index] and lines[index].strip():
        if index != start + 1:
            rows.append(
                [
                    _strip_inline_markdown(cell.strip())
                    for cell in lines[index].strip().strip("|").split("|")
                ]
            )
        index += 1
    return rows, index


def _build_docx(content: str) -> bytes:
    from docx import Document
    from docx.shared import Pt

    document = Document()
    normal_style = document.styles["Normal"]
    normal_style.font.name = "Arial"
    normal_style.font.size = Pt(11)

    lines = content.splitlines()
    index = 0
    in_code = False
    while index < len(lines):
        line = lines[index]
        stripped = line.strip()

        if stripped.startswith("```"):
            in_code = not in_code
            index += 1
            continue

        table_data = None if in_code else _markdown_rows(lines, index)
        if table_data:
            rows, index = table_data
            if rows:
                width = max(len(row) for row in rows)
                table = document.add_table(rows=len(rows), cols=width)
                table.style = "Table Grid"
                for row_index, row in enumerate(rows):
                    for column_index, value in enumerate(row):
                        table.cell(row_index, column_index).text = value
            continue

        if not stripped:
            document.add_paragraph()
        elif in_code:
            paragraph = document.add_paragraph(style="No Spacing")
            run = paragraph.add_run(line)
            run.font.name = "Consolas"
            run.font.size = Pt(9)
        elif match := re.match(r"^(#{1,6})\s+(.+)$", stripped):
            document.add_heading(
                _strip_inline_markdown(match.group(2)), level=min(len(match.group(1)), 6)
            )
        elif re.match(r"^[-*+]\s+", stripped):
            document.add_paragraph(
                _strip_inline_markdown(re.sub(r"^[-*+]\s+", "", stripped)),
                style="List Bullet",
            )
        elif re.match(r"^\d+[.)]\s+", stripped):
            document.add_paragraph(
                _strip_inline_markdown(re.sub(r"^\d+[.)]\s+", "", stripped)),
                style="List Number",
            )
        else:
            document.add_paragraph(_strip_inline_markdown(line))
        index += 1

    output = io.BytesIO()
    document.save(output)
    return output.getvalue()


def _find_unicode_font() -> str | None:
    candidates = (
        Path("C:/Windows/Fonts/arial.ttf"),
        Path("C:/Windows/Fonts/segoeui.ttf"),
        Path("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"),
        Path("/System/Library/Fonts/Supplemental/Arial Unicode.ttf"),
        Path("/Library/Fonts/Arial Unicode.ttf"),
    )
    return str(next((path for path in candidates if path.is_file()), "")) or None


def _build_pdf(content: str) -> bytes:
    from fpdf import FPDF
    from fpdf.enums import XPos, YPos

    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()

    font_path = _find_unicode_font()
    if font_path:
        pdf.add_font("NeveUnicode", fname=font_path)
        family = "NeveUnicode"
        safe_text = lambda value: value
    else:
        family = "Helvetica"
        safe_text = lambda value: value.encode("latin-1", "replace").decode("latin-1")

    in_code = False
    for line in content.splitlines() or [""]:
        stripped = line.strip()
        if stripped.startswith("```"):
            in_code = not in_code
            continue

        heading = None if in_code else re.match(r"^(#{1,6})\s+(.+)$", stripped)
        if heading:
            size = max(12, 22 - (len(heading.group(1)) * 2))
            pdf.set_font(family, size=size)
            text = _strip_inline_markdown(heading.group(2))
            line_height = max(6, size * 0.45)
        else:
            pdf.set_font(family, size=9 if in_code else 11)
            text = line if in_code else _strip_inline_markdown(line)
            if re.match(r"^[-*+]\s+", stripped):
                text = "• " + re.sub(r"^[-*+]\s+", "", stripped)
            line_height = 5 if in_code else 6

        pdf.multi_cell(
            0,
            line_height,
            safe_text(text or " "),
            new_x=XPos.LMARGIN,
            new_y=YPos.NEXT,
        )

    output = pdf.output()
    return bytes(output)


def _normalize_sheet_rows(value: Any) -> list[list[Any]]:
    if isinstance(value, dict):
        value = value.get("rows", value.get("data", []))
    if not isinstance(value, list):
        raise GeneratedFileError("Cada planilha deve conter uma lista em 'rows'.")
    if value and all(isinstance(row, dict) for row in value):
        headers = list(dict.fromkeys(key for row in value for key in row.keys()))
        return [headers, *[[row.get(header, "") for header in headers] for row in value]]
    return [row if isinstance(row, list) else [row] for row in value]


def _build_xlsx(content: str) -> bytes:
    from openpyxl import Workbook

    try:
        payload = json.loads(content)
    except json.JSONDecodeError:
        payload = {"sheets": [{"name": "Planilha", "rows": list(csv.reader(io.StringIO(content)))}]}

    if isinstance(payload, dict) and isinstance(payload.get("sheets"), list):
        sheets = payload["sheets"]
    else:
        sheets = [{"name": "Planilha", "rows": payload}]
    if not sheets:
        sheets = [{"name": "Planilha", "rows": []}]

    workbook = Workbook()
    workbook.remove(workbook.active)
    used_names: set[str] = set()

    for sheet_index, sheet in enumerate(sheets[:50], start=1):
        sheet = sheet if isinstance(sheet, dict) else {"rows": sheet}
        base_name = re.sub(r"[\\/*?:\[\]]", "_", str(sheet.get("name") or f"Planilha {sheet_index}"))[:31]
        base_name = base_name or f"Planilha {sheet_index}"
        name = base_name
        suffix = 2
        while name in used_names:
            name = f"{base_name[:27]} {suffix}"
            suffix += 1
        used_names.add(name)

        worksheet = workbook.create_sheet(name)
        rows = _normalize_sheet_rows(sheet)
        for row in rows:
            worksheet.append(row)
        if rows:
            worksheet.freeze_panes = "A2"
        for column in worksheet.columns:
            letter = column[0].column_letter
            width = min(50, max(10, max((len(str(cell.value or "")) for cell in column[:200]), default=0) + 2))
            worksheet.column_dimensions[letter].width = width

    output = io.BytesIO()
    workbook.save(output)
    return output.getvalue()


def _build_pptx(content: str) -> bytes:
    from pptx import Presentation

    try:
        payload = json.loads(content)
    except json.JSONDecodeError:
        payload = {
            "slides": [
                {"title": "", "content": section.strip()}
                for section in re.split(r"\n\s*---\s*\n", content)
                if section.strip()
            ]
        }

    slides = payload.get("slides", []) if isinstance(payload, dict) else payload
    if not isinstance(slides, list):
        raise GeneratedFileError("A apresentação deve conter uma lista em 'slides'.")

    presentation = Presentation()
    for slide_data in slides[:100]:
        slide_data = slide_data if isinstance(slide_data, dict) else {"content": slide_data}
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = str(slide_data.get("title", ""))
        body = slide.placeholders[1].text_frame
        body.clear()
        content_value = slide_data.get("content", slide_data.get("bullets", []))
        items = content_value if isinstance(content_value, list) else str(content_value).splitlines()
        for item_index, item in enumerate(items or [""]):
            paragraph = body.paragraphs[0] if item_index == 0 else body.add_paragraph()
            paragraph.text = _strip_inline_markdown(str(item))

    if not presentation.slides:
        presentation.slides.add_slide(presentation.slide_layouts[6])

    output = io.BytesIO()
    presentation.save(output)
    return output.getvalue()


def _build_zip(content: str) -> bytes:
    payload = _load_json(content, "ZIP")
    files = payload.get("files", []) if isinstance(payload, dict) else payload
    if not isinstance(files, list) or not files:
        raise GeneratedFileError("O ZIP deve conter uma lista não vazia em 'files'.")
    if len(files) > MAX_ARCHIVE_FILES:
        raise GeneratedFileError(f"O ZIP aceita no máximo {MAX_ARCHIVE_FILES} arquivos.")

    output = io.BytesIO()
    with zipfile.ZipFile(output, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        for item in files:
            if not isinstance(item, dict):
                raise GeneratedFileError("Cada item de 'files' deve ter 'path' e 'content'.")
            path = PurePosixPath(str(item.get("path", "")).replace("\\", "/"))
            if not path.name or path.is_absolute() or ".." in path.parts:
                raise GeneratedFileError(f"Caminho inválido no ZIP: {path}.")
            archive.writestr(str(path), str(item.get("content", "")).encode("utf-8"))
    return output.getvalue()


def build_generated_file(
    filename: str, content: str, requested_format: str = ""
) -> tuple[str, bytes, str]:
    if not isinstance(content, str):
        raise GeneratedFileError("O conteúdo do arquivo deve ser texto.")
    if len(content) > MAX_SOURCE_CHARS:
        raise GeneratedFileError(
            f"O conteúdo excede o limite de {MAX_SOURCE_CHARS:,} caracteres."
        )

    safe_name, file_format = _sanitize_filename(filename, requested_format)

    if file_format == "docx":
        data = _build_docx(content)
    elif file_format == "xlsx":
        data = _build_xlsx(content)
    elif file_format == "pdf":
        data = _build_pdf(content)
    elif file_format == "pptx":
        data = _build_pptx(content)
    elif file_format == "zip":
        data = _build_zip(content)
    elif file_format == "json":
        payload = _load_json(content, "JSON")
        data = json.dumps(payload, ensure_ascii=False, indent=2).encode("utf-8")
    elif file_format == "csv":
        data = content.encode("utf-8-sig")
    else:
        data = content.encode("utf-8")

    if len(data) > MAX_OUTPUT_BYTES:
        raise GeneratedFileError(
            f"O arquivo excede o limite de {MAX_OUTPUT_BYTES // (1024 * 1024)} MB."
        )

    return safe_name, data, MIME_TYPES[file_format]
